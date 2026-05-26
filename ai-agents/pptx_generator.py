import os
import sys
import glob
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Fix Windows console encoding issues for Unicode outputs
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
if hasattr(sys.stderr, "reconfigure"):
    sys.stderr.reconfigure(encoding="utf-8")

class PPTXGenerator:
    def __init__(self):
        # Color Palette - Modern Dark Theme
        self.COLOR_BG = RGBColor(0x11, 0x18, 0x27)       # slate-900 (deep gray-blue)
        self.COLOR_CARD = RGBColor(0x1F, 0x29, 0x37)     # slate-800 (card background)
        self.COLOR_BORDER = RGBColor(0x37, 0x41, 0x51)   # slate-700 (card border)
        self.COLOR_TEXT_MAIN = RGBColor(0xF9, 0xFA, 0xFB) # gray-50 (pure bright text)
        self.COLOR_TEXT_MUTED = RGBColor(0x9C, 0xA3, 0xAF) # gray-400 (secondary text)
        
        # Accent colors
        self.COLOR_ACCENT_PRIMARY = RGBColor(0x00, 0xE5, 0xFF) # Neon Cyan
        self.COLOR_ACCENT_SECONDARY = RGBColor(0xFF, 0x6B, 0x35) # Neon Orange / Coral
        self.COLOR_ACCENT_GOOD = RGBColor(0x10, 0xB9, 0x81) # Emerald Green for high scores
        
        self.FONT_TITLE = "Trebuchet MS"
        self.FONT_BODY = "Calibri"

    def apply_dark_background(self, slide):
        """Sets a solid dark background for the slide."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLOR_BG

    def add_header(self, slide, title_text, category_text=None):
        """Adds a standard premium header to the slide."""
        # Main Title Box
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = self.FONT_TITLE
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_ACCENT_PRIMARY
        
        # Sub-category or Breadcrumb if provided
        if category_text:
            p2 = tf.add_paragraph()
            p2.text = category_text.upper()
            p2.font.name = self.FONT_BODY
            p2.font.size = Pt(10)
            p2.font.bold = True
            p2.font.color.rgb = self.COLOR_TEXT_MUTED
            p2.space_before = Pt(4)

        # Draw a thin stylish accent line below the header
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.35), Inches(11.833), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = self.COLOR_BORDER
        line.line.fill.background()

    def create_title_slide(self, prs, title, genre, competitors):
        """Generates a high-end Keynote-style title slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
        self.apply_dark_background(slide)

        # Left Accent Border
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(2.2), Inches(0.08), Inches(3.0))
        border.fill.solid()
        border.fill.fore_color.rgb = self.COLOR_ACCENT_PRIMARY
        border.line.fill.background()

        # Title Text Block
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(11.0), Inches(2.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        # Main Title
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = self.FONT_TITLE
        p_title.font.size = Pt(44)
        p_title.font.bold = True
        p_title.font.color.rgb = self.COLOR_TEXT_MAIN
        p_title.space_after = Pt(12)
        
        # Subtitle
        p_sub = tf.add_paragraph()
        p_sub.text = f"Báo cáo Đánh giá Sản phẩm  |  Thể loại: {genre}"
        p_sub.font.name = self.FONT_BODY
        p_sub.font.size = Pt(18)
        p_sub.font.color.rgb = self.COLOR_ACCENT_SECONDARY
        p_sub.space_after = Pt(24)

        # Competitors
        if competitors:
            p_comp = tf.add_paragraph()
            p_comp.text = f"Đối thủ cạnh tranh chính:  {', '.join(competitors)}"
            p_comp.font.name = self.FONT_BODY
            p_comp.font.size = Pt(12)
            p_comp.font.color.rgb = self.COLOR_TEXT_MUTED

    def create_scorecard_slide(self, prs, slide_title, bullet_points):
        """Generates a premium Slide with vertical progress bars for criteria scores."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self.apply_dark_background(slide)
        self.add_header(slide, slide_title, "Đánh giá chất lượng")
        
        # Extract criteria and scores
        scores = []
        for line in bullet_points:
            match = re.search(r'-\s*\*\*(.*?)\*\*:\s*(\d+(?:\.\d+)?)\s*/\s*5\.0', line)
            if match:
                criteria = match.group(1).strip()
                score = float(match.group(2))
                scores.append((criteria, score))

        if not scores:
            # Fallback to normal text list if no match found
            self.create_bullet_slide(prs, slide_title, bullet_points)
            return

        # Scorecard Layout Calculation
        total_items = len(scores)
        y_start = Inches(1.8)
        max_height = Inches(5.0)
        
        # Limit rows to 10 max per slide to fit, split if needed
        # We cap visual bars here and scale spacing accordingly
        spacing = Inches(0.48) if total_items > 7 else Inches(0.65)
        bar_height = Inches(0.12)
        
        for idx, (criteria, score) in enumerate(scores[:10]):
            y_pos = y_start + idx * spacing
            
            # 1. Text Label (Criteria name + Score)
            lbl_box = slide.shapes.add_textbox(Inches(0.75), y_pos - Inches(0.1), Inches(5.5), Inches(0.45))
            tf = lbl_box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            
            p = tf.paragraphs[0]
            p.text = criteria
            p.font.name = self.FONT_BODY
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = self.COLOR_TEXT_MAIN
            
            p_score = p.add_run()
            p_score.text = f"  ({score}/5.0)"
            p_score.font.name = self.FONT_BODY
            p_score.font.size = Pt(13)
            p_score.font.bold = True
            # Color score differently based on value
            if score >= 4.0:
                p_score.font.color.rgb = self.COLOR_ACCENT_GOOD
            elif score >= 3.0:
                p_score.font.color.rgb = self.COLOR_ACCENT_PRIMARY
            else:
                p_score.font.color.rgb = self.COLOR_ACCENT_SECONDARY

            # 2. Progress Bar Background (gray track)
            bg_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), y_pos + Inches(0.06), Inches(6.0), bar_height)
            bg_bar.fill.solid()
            bg_bar.fill.fore_color.rgb = self.COLOR_BORDER
            bg_bar.line.fill.background()

            # 3. Progress Bar Fill
            fill_width = Inches(6.0) * (score / 5.0)
            if fill_width > 0:
                fill_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), y_pos + Inches(0.06), fill_width, bar_height)
                fill_bar.fill.solid()
                
                # Pick bar color based on score value
                if score >= 4.0:
                    fill_color = self.COLOR_ACCENT_GOOD
                elif score >= 3.0:
                    fill_color = self.COLOR_ACCENT_PRIMARY
                else:
                    fill_color = self.COLOR_ACCENT_SECONDARY
                fill_bar.fill.fore_color.rgb = fill_color
                fill_bar.line.fill.background()

    def create_column_grid_slide(self, prs, slide_title, sections, category_text="Chi tiết phân tích"):
        """Generates a Slide using card container panels for column-based layouts."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self.apply_dark_background(slide)
        self.add_header(slide, slide_title, category_text)

        num_sections = len(sections)
        if num_sections == 0:
            return

        # Width calculations based on number of columns
        total_width = Inches(11.833)
        y_pos = Inches(1.8)
        height = Inches(4.9)
        gap = Inches(0.4)
        
        # Max out at 3 columns for visual flow. If >3, render as a 2-col wrap or bullet list.
        cols = min(3, num_sections)
        width = (total_width - (gap * (cols - 1))) / cols

        for idx, (title, content) in enumerate(sections[:cols]):
            x_pos = Inches(0.75) + idx * (width + gap)

            # 1. Background Card Shape
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, y_pos, width, height)
            card.fill.solid()
            card.fill.fore_color.rgb = self.COLOR_CARD
            card.line.color.rgb = self.COLOR_BORDER
            card.line.width = Pt(1.5)

            # 2. Text Frame within the Card
            tb = slide.shapes.add_textbox(x_pos + Inches(0.2), y_pos + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

            # Section Title
            p_title = tf.paragraphs[0]
            p_title.text = title
            p_title.font.name = self.FONT_TITLE
            p_title.font.size = Pt(16)
            p_title.font.bold = True
            # Alternating header accent colors
            p_title.font.color.rgb = self.COLOR_ACCENT_PRIMARY if idx % 2 == 0 else self.COLOR_ACCENT_SECONDARY
            p_title.space_after = Pt(14)

            # Bullet points or paragraph split
            lines = [ln.strip() for ln in content.split("\n") if ln.strip()]
            for line in lines:
                p_bullet = tf.add_paragraph()
                p_bullet.font.name = self.FONT_BODY
                p_bullet.font.size = Pt(12)
                p_bullet.font.color.rgb = self.COLOR_TEXT_MAIN
                p_bullet.space_after = Pt(6)
                
                # Check for bullet signifiers
                if line.startswith("- ") or line.startswith("* ") or line.startswith("• "):
                    p_bullet.text = "• " + line[2:].strip()
                else:
                    p_bullet.text = line

    def create_bullet_slide(self, prs, slide_title, bullet_points, category_text="Nghiên cứu & Hành động"):
        """Generates a standard clean slide with a large structured bullet list."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self.apply_dark_background(slide)
        self.add_header(slide, slide_title, category_text)

        # Content Box
        tb = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.8))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

        first = True
        for line in bullet_points:
            clean_line = line.strip()
            if not clean_line:
                continue

            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()

            p.font.name = self.FONT_BODY
            p.font.size = Pt(15)
            p.font.color.rgb = self.COLOR_TEXT_MAIN
            p.space_after = Pt(12)

            # Check and format headers vs lists vs sub-items
            if clean_line.startswith("### "):
                p.text = clean_line[4:].strip()
                p.font.name = self.FONT_TITLE
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = self.COLOR_ACCENT_PRIMARY
                p.space_before = Pt(8)
                p.space_after = Pt(8)
            elif clean_line.startswith("## "):
                p.text = clean_line[3:].strip()
                p.font.name = self.FONT_TITLE
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = self.COLOR_ACCENT_SECONDARY
                p.space_before = Pt(12)
            elif clean_line.startswith("- ") or clean_line.startswith("* ") or clean_line.startswith("• "):
                p.text = "• " + clean_line[2:].strip()
                p.space_after = Pt(8)
            elif re.match(r'^\d+\.\s*', clean_line):
                # Numbered item (e.g. 1. Action item)
                p.text = clean_line
                p.font.bold = True
                p.font.color.rgb = self.COLOR_TEXT_MAIN
                p.space_after = Pt(10)
            else:
                p.text = clean_line

    def parse_slide_content(self, slide_text):
        """Helper to split slide content by sub-headings ('### 📍' or '### ')."""
        parts = re.split(r'###\s*(?:📍)?\s*(.*?)\n', slide_text)
        intro = parts[0].strip()
        sections = []
        
        for i in range(1, len(parts), 2):
            title = parts[i].strip()
            body = parts[i+1].strip() if i+1 < len(parts) else ""
            sections.append((title, body))
            
        return intro, sections

    def generate_pptx_from_markdown(self, md_path, output_dir):
        """Parses a Markdown report and outputs a beautiful presentation deck."""
        if not os.path.exists(md_path):
            print(f"Markdown file not found: {md_path}")
            return None
            
        filename = os.path.basename(md_path)
        print(f"Reading and parsing: {filename}")
        
        with open(md_path, "r", encoding="utf-8") as f:
            content = f.read()

        # Split content into slides using "---" divider
        slide_raw_blocks = content.split("---")
        if not slide_raw_blocks:
            print("No slides found in Markdown.")
            return None

        # 1. Parse Project Metadata from Title & Slide 1
        project_name = "Unknown Project"
        genre = "Casual"
        competitors = []
        
        # Parse main header
        main_title_match = re.search(r'^#\s*Báo cáo Đánh giá sản phẩm:\s*(.*?)\n', content)
        if main_title_match:
            project_name = main_title_match.group(1).strip()
            
        # Parse Slide 1 block details
        slide1_text = slide_raw_blocks[0]
        for line in slide1_text.split("\n"):
            if "thể loại" in line.lower() or "genre" in line.lower():
                val = line.split(":")[-1].strip()
                if val: genre = val
            elif "đối thủ" in line.lower() or "competitors" in line.lower():
                val = line.split(":")[-1].strip()
                if val and val != "Không có (N/A)":
                    competitors = [c.strip() for c in val.split(",") if c.strip()]

        # Initialize python-pptx presentation
        prs = Presentation()
        prs.slide_width = Inches(13.333) # 16:9 widescreen width
        prs.slide_height = Inches(7.5)   # 16:9 widescreen height

        # Create Slide 1 (Title Slide)
        self.create_title_slide(prs, project_name, genre, competitors)

        # Parse subsequent slides
        for block in slide_raw_blocks[1:]:
            block_lines = [ln.strip() for ln in block.split("\n") if ln.strip()]
            if not block_lines:
                continue

            # Detect slide header
            slide_header = "Slide Content"
            header_idx = -1
            for idx, line in enumerate(block_lines):
                if line.startswith("## Slide"):
                    slide_header = re.sub(r'## Slide \d+:\s*', '', line).strip()
                    header_idx = idx
                    break

            # Content body below header
            content_lines = block_lines[header_idx+1:] if header_idx != -1 else block_lines
            content_text = "\n".join(content_lines)

            # Slide 2: Scorecard check
            if "scorecard" in slide_header.lower() or "bảng điểm" in slide_header.lower():
                self.create_scorecard_slide(prs, slide_header, content_lines)
                continue

            # Standard layout parsing (Column cards vs bullets)
            intro, sections = self.parse_slide_content(content_text)
            
            if sections:
                # If we have distinct 📍 sections, render them in a column grid layout
                category = "Phân tích chi tiết"
                if "swot" in slide_header.lower():
                    category = "Ma trận SWOT"
                elif "hành động" in slide_header.lower() or "hành vi" in slide_header.lower():
                    category = "Kế hoạch hành động"
                
                self.create_column_grid_slide(prs, slide_header, sections, category)
            else:
                # Fallback to standard structured list slide
                category = "Tổng kết thông tin"
                if "nghiên cứu" in slide_header.lower() or "đối thủ" in slide_header.lower():
                    category = "Nghiên cứu đối thủ"
                elif "hành động" in slide_header.lower() or "action" in slide_header.lower():
                    category = "Đề xuất hành động"
                self.create_bullet_slide(prs, slide_header, content_lines, category)

        # Save Presentation
        safe_name = re.sub(r'[\\/*?:"<>|]', "", project_name).replace(" ", "_")
        output_path = os.path.join(output_dir, f"Presentation_{safe_name}.pptx")
        prs.save(output_path)
        print(f"Presentation successfully saved at: {output_path}")
        return output_path

def main():
    import argparse
    parser = argparse.ArgumentParser(description="PPTX Slide Deck Compiler")
    parser.add_argument("--markdown", type=str, help="Specific markdown file path to compile")
    args = parser.parse_args()

    workspace_dir = "c:/Users/Admin/Desktop/agent AI"
    output_dir = os.path.join(workspace_dir, "workflow-nhan/PowerPoint")
    os.makedirs(output_dir, exist_ok=True)

    print("=== PPTX Slide Deck Compiler ===")
    
    compiler = PPTXGenerator()

    if args.markdown:
        md_path = args.markdown
        if not os.path.isabs(md_path):
            md_path = os.path.abspath(md_path)
        print(f"Compiling single file: {md_path}")
        try:
            res = compiler.generate_pptx_from_markdown(md_path, output_dir)
            if res:
                print(f"\nCompleted! Presentation saved at: {res}")
            else:
                print("\nCompilation failed.")
        except Exception as e:
            print(f"Error compiling {os.path.basename(md_path)}: {e}")
            import traceback
            traceback.print_exc()
    else:
        md_dir = os.path.join(workspace_dir, "workflow-nhan/Đánh giá của AI")
        print(f"Scanning directory: {md_dir}")
        print(f"Output directory: {output_dir}")

        md_files = glob.glob(os.path.join(md_dir, "*.md"))
        if not md_files:
            print("No evaluation Markdown reports found.")
            return

        print(f"Found {len(md_files)} Markdown report(s) to convert.")
        compiled_count = 0
        for md_path in md_files:
            try:
                res = compiler.generate_pptx_from_markdown(md_path, output_dir)
                if res:
                    compiled_count += 1
            except Exception as e:
                print(f"Error compiling {os.path.basename(md_path)}: {e}")
                import traceback
                traceback.print_exc()

        print(f"\nCompleted! Successfully compiled {compiled_count}/{len(md_files)} presentations.")

if __name__ == "__main__":
    main()
